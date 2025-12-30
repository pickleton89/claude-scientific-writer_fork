"""
Oligon Brand Adapter for python-pptx (PowerPoint Presentations)

This module provides functions and styles for creating PowerPoint presentations
with Oligon scientific brand identity using python-pptx.

Usage:
    from oligon_brand.adapters.pptx_adapter import (
        create_brand_presentation,
        add_title_slide,
        add_content_slide,
        BRAND_COLORS,
    )

    prs = create_brand_presentation()
    add_title_slide(prs, "Presentation Title", "Subtitle here")
    add_content_slide(prs, "Section Title", ["Bullet 1", "Bullet 2"])
    prs.save("output.pptx")
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from typing import Optional, List, Tuple

# =============================================================================
# BRAND COLORS
# =============================================================================

BRAND_COLORS = {
    # Primary highlight
    'brand_blue': RGBColor(0x2D, 0xB2, 0xE8),

    # Neutrals
    'dark_gray': RGBColor(0x22, 0x22, 0x22),
    'medium_gray': RGBColor(0x66, 0x66, 0x66),
    'muted_gray': RGBColor(0x99, 0x99, 0x99),
    'light_gray': RGBColor(0xBD, 0xBD, 0xBD),

    # Contrast
    'contrast_orange': RGBColor(0xE8, 0x62, 0x2D),
    'dark_warm': RGBColor(0x7D, 0x25, 0x0F),

    # Supporting blues
    'medium_blue': RGBColor(0x15, 0x8B, 0xBB),
    'dark_teal': RGBColor(0x0F, 0x5D, 0x7D),

    # Structure
    'black': RGBColor(0x00, 0x00, 0x00),
    'white': RGBColor(0xFF, 0xFF, 0xFF),
}

# =============================================================================
# TYPOGRAPHY SETTINGS
# =============================================================================

FONT_NAME = 'Arial'

FONT_SIZES = {
    'title': Pt(28),
    'subtitle': Pt(20),
    'heading': Pt(24),
    'body': Pt(18),
    'bullet': Pt(16),
    'caption': Pt(12),
    'footer': Pt(10),
}

# =============================================================================
# SLIDE DIMENSIONS (16:9 widescreen)
# =============================================================================

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Content margins
MARGIN_LEFT = Inches(0.5)
MARGIN_RIGHT = Inches(0.5)
MARGIN_TOP = Inches(0.5)
MARGIN_BOTTOM = Inches(0.5)

# =============================================================================
# PRESENTATION CREATION
# =============================================================================

def create_brand_presentation() -> Presentation:
    """
    Create a new PowerPoint presentation with brand settings.

    Returns
    -------
    Presentation
        A python-pptx Presentation configured for brand compliance

    Examples
    --------
    >>> prs = create_brand_presentation()
    >>> add_title_slide(prs, "My Presentation")
    >>> prs.save("presentation.pptx")
    """
    prs = Presentation()

    # Set slide dimensions (16:9)
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    return prs


# =============================================================================
# SLIDE CREATION FUNCTIONS
# =============================================================================

def add_title_slide(
    prs: Presentation,
    title: str,
    subtitle: Optional[str] = None
) -> 'Slide':
    """
    Add a title slide with brand styling.

    Parameters
    ----------
    prs : Presentation
        The presentation
    title : str
        Main title text
    subtitle : str, optional
        Subtitle text

    Returns
    -------
    Slide
        The created slide
    """
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Set white background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BRAND_COLORS['white']

    # Add title
    title_box = slide.shapes.add_textbox(
        MARGIN_LEFT,
        Inches(2.5),
        SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT,
        Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.name = FONT_NAME
    title_para.font.size = FONT_SIZES['title']
    title_para.font.bold = True
    title_para.font.color.rgb = BRAND_COLORS['black']
    title_para.alignment = PP_ALIGN.CENTER

    # Add subtitle if provided
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(
            MARGIN_LEFT,
            Inches(4.2),
            SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT,
            Inches(1)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.text = subtitle
        subtitle_para.font.name = FONT_NAME
        subtitle_para.font.size = FONT_SIZES['subtitle']
        subtitle_para.font.color.rgb = BRAND_COLORS['medium_gray']
        subtitle_para.alignment = PP_ALIGN.CENTER

    # Add brand accent line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(4),
        Inches(4),
        Inches(5.333),
        Pt(4)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = BRAND_COLORS['brand_blue']
    line.line.fill.background()

    return slide


def add_content_slide(
    prs: Presentation,
    title: str,
    bullets: Optional[List[str]] = None,
    content_text: Optional[str] = None
) -> 'Slide':
    """
    Add a content slide with title and bullets or text.

    Parameters
    ----------
    prs : Presentation
        The presentation
    title : str
        Slide title
    bullets : list, optional
        List of bullet points
    content_text : str, optional
        Paragraph text (alternative to bullets)

    Returns
    -------
    Slide
        The created slide
    """
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Set white background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BRAND_COLORS['white']

    # Add title
    title_box = slide.shapes.add_textbox(
        MARGIN_LEFT,
        MARGIN_TOP,
        SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT,
        Inches(1)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.name = FONT_NAME
    title_para.font.size = FONT_SIZES['heading']
    title_para.font.bold = True
    title_para.font.color.rgb = BRAND_COLORS['black']

    # Add content area
    content_top = Inches(1.5)
    content_height = SLIDE_HEIGHT - content_top - MARGIN_BOTTOM

    if bullets:
        content_box = slide.shapes.add_textbox(
            MARGIN_LEFT,
            content_top,
            SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT,
            content_height
        )
        text_frame = content_box.text_frame
        text_frame.word_wrap = True

        for i, bullet in enumerate(bullets):
            if i == 0:
                para = text_frame.paragraphs[0]
            else:
                para = text_frame.add_paragraph()

            para.text = bullet
            para.font.name = FONT_NAME
            para.font.size = FONT_SIZES['bullet']
            para.font.color.rgb = BRAND_COLORS['dark_gray']
            para.level = 0
            para.space_after = Pt(12)

    elif content_text:
        content_box = slide.shapes.add_textbox(
            MARGIN_LEFT,
            content_top,
            SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT,
            content_height
        )
        text_frame = content_box.text_frame
        text_frame.word_wrap = True
        para = text_frame.paragraphs[0]
        para.text = content_text
        para.font.name = FONT_NAME
        para.font.size = FONT_SIZES['body']
        para.font.color.rgb = BRAND_COLORS['dark_gray']

    return slide


def add_section_slide(
    prs: Presentation,
    title: str
) -> 'Slide':
    """
    Add a section divider slide with brand styling.

    Parameters
    ----------
    prs : Presentation
        The presentation
    title : str
        Section title

    Returns
    -------
    Slide
        The created slide
    """
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Brand blue background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BRAND_COLORS['brand_blue']

    # White title text
    title_box = slide.shapes.add_textbox(
        MARGIN_LEFT,
        Inches(3),
        SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT,
        Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.name = FONT_NAME
    title_para.font.size = FONT_SIZES['title']
    title_para.font.bold = True
    title_para.font.color.rgb = BRAND_COLORS['white']
    title_para.alignment = PP_ALIGN.CENTER

    return slide


def add_two_column_slide(
    prs: Presentation,
    title: str,
    left_content: List[str],
    right_content: List[str]
) -> 'Slide':
    """
    Add a two-column content slide.

    Parameters
    ----------
    prs : Presentation
        The presentation
    title : str
        Slide title
    left_content : list
        Bullet points for left column
    right_content : list
        Bullet points for right column

    Returns
    -------
    Slide
        The created slide
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # White background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BRAND_COLORS['white']

    # Title
    title_box = slide.shapes.add_textbox(
        MARGIN_LEFT,
        MARGIN_TOP,
        SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT,
        Inches(1)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.name = FONT_NAME
    title_para.font.size = FONT_SIZES['heading']
    title_para.font.bold = True
    title_para.font.color.rgb = BRAND_COLORS['black']

    # Calculate column dimensions
    content_top = Inches(1.5)
    content_height = SLIDE_HEIGHT - content_top - MARGIN_BOTTOM
    column_width = (SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT - Inches(0.5)) / 2

    # Left column
    left_box = slide.shapes.add_textbox(
        MARGIN_LEFT,
        content_top,
        column_width,
        content_height
    )
    _add_bullets_to_textbox(left_box, left_content)

    # Right column
    right_box = slide.shapes.add_textbox(
        MARGIN_LEFT + column_width + Inches(0.5),
        content_top,
        column_width,
        content_height
    )
    _add_bullets_to_textbox(right_box, right_content)

    return slide


def _add_bullets_to_textbox(textbox, bullets: List[str]) -> None:
    """Helper to add bullet points to a textbox."""
    text_frame = textbox.text_frame
    text_frame.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            para = text_frame.paragraphs[0]
        else:
            para = text_frame.add_paragraph()

        para.text = bullet
        para.font.name = FONT_NAME
        para.font.size = FONT_SIZES['bullet']
        para.font.color.rgb = BRAND_COLORS['dark_gray']
        para.space_after = Pt(8)


# =============================================================================
# HELPER FUNCTIONS
# =============================================================================

def add_image_slide(
    prs: Presentation,
    title: str,
    image_path: str,
    caption: Optional[str] = None
) -> 'Slide':
    """
    Add a slide with an image and optional caption.

    Parameters
    ----------
    prs : Presentation
        The presentation
    title : str
        Slide title
    image_path : str
        Path to the image file
    caption : str, optional
        Image caption

    Returns
    -------
    Slide
        The created slide
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # White background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BRAND_COLORS['white']

    # Title
    title_box = slide.shapes.add_textbox(
        MARGIN_LEFT,
        MARGIN_TOP,
        SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT,
        Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.name = FONT_NAME
    title_para.font.size = FONT_SIZES['heading']
    title_para.font.bold = True
    title_para.font.color.rgb = BRAND_COLORS['black']

    # Image (centered)
    img_top = Inches(1.3)
    img_height = Inches(5) if caption else Inches(5.5)

    # Add image centered
    slide.shapes.add_picture(
        image_path,
        Inches(1.5),
        img_top,
        height=img_height
    )

    # Caption
    if caption:
        caption_box = slide.shapes.add_textbox(
            MARGIN_LEFT,
            Inches(6.5),
            SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT,
            Inches(0.5)
        )
        caption_frame = caption_box.text_frame
        caption_para = caption_frame.paragraphs[0]
        caption_para.text = caption
        caption_para.font.name = FONT_NAME
        caption_para.font.size = FONT_SIZES['caption']
        caption_para.font.italic = True
        caption_para.font.color.rgb = BRAND_COLORS['medium_gray']
        caption_para.alignment = PP_ALIGN.CENTER

    return slide


def add_highlight_box(
    slide,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    bg_color: str = 'brand_blue',
    text_color: str = 'white'
) -> None:
    """
    Add a highlighted text box to a slide.

    Parameters
    ----------
    slide : Slide
        The slide to add the box to
    text : str
        Text content
    left, top, width, height : float
        Position and size in inches
    bg_color : str
        Background color name
    text_color : str
        Text color name
    """
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height)
    )

    shape.fill.solid()
    shape.fill.fore_color.rgb = BRAND_COLORS.get(bg_color, BRAND_COLORS['brand_blue'])
    shape.line.fill.background()

    text_frame = shape.text_frame
    text_frame.word_wrap = True
    para = text_frame.paragraphs[0]
    para.text = text
    para.font.name = FONT_NAME
    para.font.size = FONT_SIZES['body']
    para.font.bold = True
    para.font.color.rgb = BRAND_COLORS.get(text_color, BRAND_COLORS['white'])
    para.alignment = PP_ALIGN.CENTER
    text_frame.anchor = MSO_ANCHOR.MIDDLE


# =============================================================================
# EXAMPLE USAGE
# =============================================================================

def create_example_presentation():
    """Create an example presentation demonstrating brand styles."""
    prs = create_brand_presentation()

    # Title slide
    add_title_slide(
        prs,
        "Oligon Scientific Presentation",
        "Demonstrating Brand-Compliant Slides"
    )

    # Section divider
    add_section_slide(prs, "Introduction")

    # Content slide with bullets
    add_content_slide(
        prs,
        "Key Findings",
        bullets=[
            "Treatment group showed 62% improvement",
            "Statistical significance achieved (p < 0.001)",
            "Effect size was large (Cohen's d = 1.2)",
            "No adverse events reported"
        ]
    )

    # Two-column slide
    add_two_column_slide(
        prs,
        "Comparison",
        left_content=[
            "Control Group",
            "N = 25 participants",
            "Mean response: 4.2",
            "SD: 1.1"
        ],
        right_content=[
            "Treatment Group",
            "N = 25 participants",
            "Mean response: 6.8",
            "SD: 0.9"
        ]
    )

    # Section divider
    add_section_slide(prs, "Conclusions")

    # Final content slide
    add_content_slide(
        prs,
        "Summary",
        content_text=(
            "This study demonstrates significant efficacy of the treatment "
            "approach. Results support further investigation in larger trials."
        )
    )

    prs.save("example_brand_presentation.pptx")
    print("Example presentation created: example_brand_presentation.pptx")


if __name__ == "__main__":
    create_example_presentation()
